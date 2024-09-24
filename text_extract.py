import os
import PyPDF2
import docx
from pptx import Presentation
import re
import logging
from utils import count_tokens

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def clean_text(text):
    # Remove extra whitespace
    text = re.sub(r'\s+', ' ', text).strip()
    # Remove non-printable characters
    text = ''.join(char for char in text if char.isprintable() or char.isspace())
    return text

def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        logger.error(f"Error extracting text from PDF {pdf_path}: {str(e)}")
    return clean_text(text)

def extract_text_from_docx(docx_path):
    try:
        doc = docx.Document(docx_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return clean_text(text)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {docx_path}: {str(e)}")
        return ""

def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return clean_text(text)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {pptx_path}: {str(e)}")
        return ""

def extract_text_from_txt(txt_path):
    try:
        with open(txt_path, 'r', encoding='utf-8') as file:
            text = file.read()
        return clean_text(text)
    except Exception as e:
        logger.error(f"Error extracting text from TXT {txt_path}: {str(e)}")
        return ""

def extract_text_from_files(path):
    all_text = ""

    if os.path.isdir(path):
        for filename in os.listdir(path):
            file_path = os.path.join(path, filename)
            logger.info(f"Processing file: {filename}")

            if filename.lower().endswith('.pdf'):
                text = extract_text_from_pdf(file_path)
            elif filename.lower().endswith('.docx'):
                text = extract_text_from_docx(file_path)
            elif filename.lower().endswith('.pptx'):
                text = extract_text_from_pptx(file_path)
            elif filename.lower().endswith('.txt'):
                text = extract_text_from_txt(file_path)
            else:
                logger.warning(f"Unsupported file type: {filename}")
                continue

            if text:
                all_text += f"--- Start of {filename} ---\n"
                all_text += text
                all_text += f"\n--- End of {filename} ---\n\n"
            else:
                logger.warning(f"No text extracted from {filename}")

    elif os.path.isfile(path):
        filename = os.path.basename(path)
        logger.info(f"Processing file: {filename}")

        if filename.lower().endswith('.pdf'):
            text = extract_text_from_pdf(path)
        elif filename.lower().endswith('.docx'):
            text = extract_text_from_docx(path)
        elif filename.lower().endswith('.pptx'):
            text = extract_text_from_pptx(path)
        elif filename.lower().endswith('.txt'):
            text = extract_text_from_txt(path)
        else:
            logger.warning(f"Unsupported file type: {filename}")
            text = None  # Set text to None for unsupported file types

        if text:
            all_text += f"--- Start of {filename} ---\n"
            all_text += text
            all_text += f"\n--- End of {filename} ---\n\n"
        else:
            logger.warning(f"No text extracted from {filename}")

    else:
        logger.error(f"Invalid path: {path}")

    return all_text